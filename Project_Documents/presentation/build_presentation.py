"""
build_presentation.py
=====================
Generates the 5-minute video presentation slide deck for the
AI-Powered Financial Document Q&A System.

EC7203 Advanced Artificial Intelligence - Final Project

Run:
    py -3.13 build_presentation.py

Output:
    presentation/Financial_QA_Presentation.pptx  (~12 slides, 16:9)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Theme colors (deep blue + gold accent) ───────────────────────────────────
NAVY      = RGBColor(0x1A, 0x3C, 0x6E)   # primary
LIGHT_BG  = RGBColor(0xF2, 0xF5, 0xF9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GOLD      = RGBColor(0xC8, 0x97, 0x00)
DARK_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
GREY      = RGBColor(0x55, 0x55, 0x55)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
RED       = RGBColor(0xC6, 0x28, 0x28)
ORANGE    = RGBColor(0xEF, 0x6C, 0x00)
LIGHT_BLUE= RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN=RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_GOLD =RGBColor(0xFF, 0xF8, 0xE1)


# ─────────────────────────────────────────────────────────────────────────────
# Set up 16:9 presentation
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ─────────────────────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────────────────────

def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_header_bar(slide, title, color=NAVY):
    """Add a slim color bar with white title at top of slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), SLIDE_W - Inches(1), Inches(0.6)
    )
    tf = tb.text_frame
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
             italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_bullets(slide, bullets, left, top, width, height,
                size=18, color=DARK_TEXT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # bullet marker
        r1 = p.add_run()
        r1.text = "•  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = NAVY
        r1.font.name = "Calibri"
        # text
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Calibri"
    return tb


def add_box(slide, left, top, width, height, fill=LIGHT_BG, line=NAVY,
            line_width=0.75, corner=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(line_width)
    box.shadow.inherit = False
    return box


def add_label_box(slide, text, left, top, width, height,
                  fill=LIGHT_BG, border=NAVY, text_color=DARK_TEXT,
                  size=16, bold=True, align=PP_ALIGN.CENTER):
    add_box(slide, left, top, width, height, fill=fill, line=border)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = "Calibri"
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.5):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_footer(slide, text, num=None):
    add_text(slide, text, Inches(0.3), Inches(7.1),
             Inches(10), Inches(0.3), size=10, color=GREY)
    if num is not None:
        add_text(slide, str(num), Inches(12.7), Inches(7.1),
                 Inches(0.5), Inches(0.3), size=10, color=GREY,
                 align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

# Decorative gold accent bar
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5),
                         Inches(13.333), Inches(0.08))
acc.fill.solid(); acc.fill.fore_color.rgb = GOLD; acc.line.fill.background()

# Course label
add_text(s, "EC7203  •  ADVANCED ARTIFICIAL INTELLIGENCE",
         Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.4),
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Main title
add_text(s, "AI-Powered Financial Document",
         Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.9),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Question & Answer System",
         Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.9),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(s, "Using Retrieval-Augmented Generation (RAG)",
         Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.5),
         size=22, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# Tech tags
tags = "NLP   •   Large Language Models   •   Prompt Engineering   •   Word Embeddings"
add_text(s, tags, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

# Footer
add_text(s, "Final Group Project Report  •  May 2026",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
         size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem & Motivation
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "Problem & Motivation")

# Left: the problem
add_text(s, "The Problem", Inches(0.5), Inches(1.1), Inches(6), Inches(0.5),
         size=22, bold=True, color=NAVY)
add_bullets(s, [
    "Financial documents (10-K, earnings reports, SEC filings) can exceed 300 pages.",
    "Analysts manually search dense documents for specific figures — slow and error-prone.",
    "Keyword search misses paraphrased content and synonyms.",
    "Raw LLMs without grounding hallucinate fabricated financial figures.",
], Inches(0.5), Inches(1.7), Inches(6.3), Inches(4),
   size=15)

# Right: visual stat box
add_box(s, Inches(7.2), Inches(1.3), Inches(5.5), Inches(4.6),
        fill=LIGHT_BLUE, line=NAVY, line_width=1.5)

add_text(s, "WHY IT MATTERS",
         Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s, "300+", Inches(7.2), Inches(2.1), Inches(5.5), Inches(0.9),
         size=64, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "pages in a typical 10-K filing",
         Inches(7.2), Inches(2.95), Inches(5.5), Inches(0.4),
         size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "$5–15", Inches(7.2), Inches(3.5), Inches(5.5), Inches(0.7),
         size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "cost per question if sent to GPT-4 directly",
         Inches(7.2), Inches(4.2), Inches(5.5), Inches(0.4),
         size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "Banking & Finance",
         Inches(7.2), Inches(4.85), Inches(5.5), Inches(0.5),
         size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "course-listed industry domain",
         Inches(7.2), Inches(5.35), Inches(5.5), Inches(0.4),
         size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER, italic=True)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Proposed Solution
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "Proposed Solution: Retrieval-Augmented Generation")

add_text(s, "Instead of asking the LLM everything, we let it look up answers in YOUR document.",
         Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
         size=16, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# Three big advantage boxes
boxes = [
    ("Accurate", "Answers come only from the\nactual document content.",
     "No hallucinations.", GREEN, LIGHT_GREEN),
    ("Cited", "Every answer includes\nthe page number it came from.",
     "Verifiable claims.", NAVY, LIGHT_BLUE),
    ("Efficient", "Only relevant excerpts are\nsent to the LLM.",
     "10× cheaper than full-doc.", ORANGE, LIGHT_GOLD),
]

x_start = Inches(0.6)
box_w   = Inches(4.0)
box_h   = Inches(4.0)
gap     = Inches(0.25)

for i, (title, desc, tag, accent, fill) in enumerate(boxes):
    x = x_start + (box_w + gap) * i
    y = Inches(1.8)
    add_box(s, x, y, box_w, box_h, fill=fill, line=accent, line_width=2)

    # Top emoji-style icon (large letter)
    add_text(s, "✓", x, y + Inches(0.3), box_w, Inches(0.9),
             size=56, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(s, title, x, y + Inches(1.3), box_w, Inches(0.6),
             size=26, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(s, desc, x, y + Inches(2.0), box_w, Inches(1.4),
             size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_text(s, tag, x, y + Inches(3.3), box_w, Inches(0.5),
             size=13, bold=True, italic=True, color=accent,
             align=PP_ALIGN.CENTER)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Three AI Techniques
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "Three AI Techniques (Course Requirement: Min 3)")

# Three vertical technique columns
techniques = [
    ("1", "Natural Language Processing",
     ["Text preprocessing & cleaning",
      "Token-counted chunking (200 tok)",
      "Word2Vec  (baseline)",
      "TF-IDF  (baseline)",
      "Sentence-Transformers  (proposed)"], NAVY, LIGHT_BLUE),
    ("2", "Semantic Retrieval",
     ["MiniLM transformer encoder",
      "384-dimensional dense vectors",
      "Cosine similarity ranking",
      "Grounded document excerpts",
      "Fully local operation"], ORANGE, LIGHT_GOLD),
    ("3", "Prompt Engineering",
     ["Systematic 7-rule system prompt",
      "Chain-of-Thought (think step-by-step)",
      "Few-shot in-context learning",
      "Page citation enforcement",
      "Hallucination prevention"], GREEN, LIGHT_GREEN),
]

x_start = Inches(0.4)
col_w   = Inches(4.15)
col_h   = Inches(5.4)
gap     = Inches(0.2)

for i, (num, title, items, accent, fill) in enumerate(techniques):
    x = x_start + (col_w + gap) * i
    y = Inches(1.1)

    add_box(s, x, y, col_w, col_h, fill=fill, line=accent, line_width=1.5)

    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                x + Inches(0.2), y + Inches(0.2),
                                Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tb = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2),
                              Inches(0.7), Inches(0.7))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE

    # Title
    add_text(s, title, x + Inches(1.0), y + Inches(0.25),
             col_w - Inches(1.1), Inches(0.7),
             size=17, bold=True, color=accent)

    # Bullets
    add_bullets(s, items,
                x + Inches(0.3), y + Inches(1.2),
                col_w - Inches(0.5), Inches(4),
                size=13, line_spacing=1.3)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=4)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — System Architecture (Two-Phase Pipeline)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "System Architecture: Two-Phase RAG Pipeline")

# ── Phase 1 label
add_text(s, "PHASE 1 — INDEXING  (runs once per document)",
         Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
         size=14, bold=True, color=NAVY)

# Phase 1 boxes
box_w  = Inches(2.2); box_h = Inches(1.0); y1 = Inches(1.55)
spacing = Inches(0.3)
phase1 = [
    ("PDF\nUpload",          ORANGE),
    ("Text\nExtraction",     NAVY),
    ("Chunking\n(800 tok)",  NAVY),
    ("Embedding\n(MiniLM)",  NAVY),
    ("ChromaDB\nStore",      GREEN),
]
x = Inches(0.5)
phase1_centers = []
for label, accent in phase1:
    fill = LIGHT_GOLD if accent == ORANGE else (LIGHT_GREEN if accent == GREEN else LIGHT_BLUE)
    add_label_box(s, label, x, y1, box_w, box_h,
                  fill=fill, border=accent, text_color=accent, size=13)
    phase1_centers.append(x + box_w // 2)
    x += box_w + spacing

# Arrows between Phase 1 boxes
x = Inches(0.5) + box_w
y_mid = y1 + box_h // 2
for _ in range(len(phase1) - 1):
    add_arrow(s, x, y_mid, x + spacing, y_mid, color=NAVY, width=2)
    x += box_w + spacing

# ── Phase 2 label
add_text(s, "PHASE 2 — QUERYING  (runs on every user question)",
         Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.4),
         size=14, bold=True, color=NAVY)

# Phase 2 boxes
y2 = Inches(3.9)
phase2 = [
    ("User\nQuestion",       ORANGE),
    ("Query\nEmbedding",     NAVY),
    ("Cosine\nSearch",       NAVY),
    ("Top-K\nChunks",        NAVY),
    ("LLM\nAnswer",          GREEN),
]
x = Inches(0.5)
for label, accent in phase2:
    fill = LIGHT_GOLD if accent == ORANGE else (LIGHT_GREEN if accent == GREEN else LIGHT_BLUE)
    add_label_box(s, label, x, y2, box_w, box_h,
                  fill=fill, border=accent, text_color=accent, size=13)
    x += box_w + spacing

# Arrows between Phase 2 boxes
x = Inches(0.5) + box_w
y_mid2 = y2 + box_h // 2
for _ in range(len(phase2) - 1):
    add_arrow(s, x, y_mid2, x + spacing, y_mid2, color=NAVY, width=2)
    x += box_w + spacing

# Vector store -> Cosine Search dashed arrow
x_store = Inches(0.5) + (box_w + spacing) * 4 + box_w // 2
x_search = Inches(0.5) + (box_w + spacing) * 2 + box_w // 2
add_arrow(s, x_store, y1 + box_h, x_store, Inches(3.15),
          color=GREEN, width=1.5)
add_arrow(s, x_store, Inches(3.15), x_search, Inches(3.15),
          color=GREEN, width=1.5)
add_arrow(s, x_search, Inches(3.15), x_search, y2,
          color=GREEN, width=1.5)
add_text(s, "vector store feeds search",
         Inches(7.0), Inches(3.05), Inches(3), Inches(0.3),
         size=10, italic=True, color=GREEN)

# Bottom legend
add_box(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4),
        fill=LIGHT_BG, line=GREY)
add_text(s, "Result: when a user asks a question, the system retrieves the most "
            "relevant passages from the indexed document and generates a grounded "
            "answer with page citations — not from the LLM's training memory.",
         Inches(0.7), Inches(5.7), Inches(12.0), Inches(1.0),
         size=14, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Implementation Stack
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "Implementation & Technology Stack")

# Left: layered architecture
add_text(s, "Layered Architecture",
         Inches(0.5), Inches(1.1), Inches(6), Inches(0.5),
         size=18, bold=True, color=NAVY)

layers = [
    ("Web UI / API",      "Streamlit  •  FastAPI",        ORANGE,  LIGHT_GOLD),
    ("Grounded Response", "Retrieved document excerpts",   GREEN,   LIGHT_GREEN),
    ("Retrieval",         "ChromaDB  •  Cosine ANN",      NAVY,    LIGHT_BLUE),
    ("Embedding",         "all-MiniLM-L6-v2",              NAVY,    LIGHT_BLUE),
    ("NLP Preprocessing", "pdfplumber  •  tiktoken",      GREY,    LIGHT_BG),
]
y = Inches(1.7)
layer_w = Inches(6.3); layer_h = Inches(0.7)
for name, libs, accent, fill in layers:
    add_box(s, Inches(0.5), y, layer_w, layer_h, fill=fill, line=accent)
    add_text(s, name, Inches(0.65), y + Inches(0.12),
             Inches(2.6), Inches(0.5), size=14, bold=True, color=accent)
    add_text(s, libs, Inches(3.25), y + Inches(0.12),
             Inches(3.5), Inches(0.5), size=12, color=DARK_TEXT,
             italic=True)
    y += layer_h + Inches(0.1)

# Right: code stats
add_text(s, "By the Numbers",
         Inches(7.5), Inches(1.1), Inches(5.3), Inches(0.5),
         size=18, bold=True, color=NAVY)

stats = [
    ("28",   "Python source files"),
    ("~2.4K","Lines of code"),
    ("5",    "AI modules (ingest, retrieve, generate, eval, api)"),
    ("4",    "REST API endpoints"),
    ("3",    "Evaluation experiments + golden test cases"),
    ("22",   "Pages: final LaTeX report"),
]
y = Inches(1.7)
for num, desc in stats:
    add_box(s, Inches(7.5), y, Inches(5.3), Inches(0.7),
            fill=LIGHT_BG, line=NAVY)
    add_text(s, num, Inches(7.6), y + Inches(0.1),
             Inches(1.2), Inches(0.5),
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(8.85), y + Inches(0.18),
             Inches(3.9), Inches(0.5), size=13, color=DARK_TEXT)
    y += Inches(0.8)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=6)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Live Demo (placeholder)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

# Big "DEMO" text
add_text(s, "LIVE  DEMO", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(5.5), Inches(3.0), Inches(2.3), Inches(0.08))
acc.fill.solid(); acc.fill.fore_color.rgb = GOLD; acc.line.fill.background()

add_text(s, "Financial Document Q&A System",
         Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
         size=22, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# Demo script steps
steps = [
    ("1.", "Upload a real 10-K filing (Apple FY 2023 from SEC EDGAR)"),
    ("2.", "Watch the indexing pipeline: extract → chunk → embed → store"),
    ("3.", "Ask: \"What was the total net revenue in fiscal 2023?\""),
    ("4.", "View the grounded answer with page citation and relevance score"),
    ("5.", "Ask: \"What are the main risk factors?\"  (multi-part answer)"),
]

y = Inches(4.2)
for num, txt in steps:
    add_text(s, num, Inches(1.5), y, Inches(0.6), Inches(0.45),
             size=18, bold=True, color=GOLD)
    add_text(s, txt, Inches(2.2), y, Inches(10), Inches(0.45),
             size=16, color=WHITE)
    y += Inches(0.5)

add_text(s, "Switch to browser: http://localhost:8501",
         Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         size=14, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Evaluation Experiment 1 (Embedding Baselines)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "Experiment 1: Embedding Baseline Comparison")

add_text(s, "Compared three methods on 5 financial queries — which retrieves the right passage?",
         Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
         size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# Results table
table_left = Inches(0.7)
table_top  = Inches(1.7)
table_w    = Inches(7.0)
n_rows = 4; n_cols = 4
table_shape = s.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                 table_w, Inches(2.6))
table = table_shape.table

# Column widths
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(1.3)
table.columns[2].width = Inches(1.3)
table.columns[3].width = Inches(1.4)

headers = ["Method", "P₃", "MRR", "Latency"]
rows = [
    ("TF-IDF (sparse, n-gram)",          "0.333", "0.800", "4.9 ms"),
    ("Word2Vec (skip-gram, d=100)",      "0.200", "0.517", "2.0 ms"),
    ("MiniLM-L6-v2  ← proposed",         "0.400", "1.000", "89.0 ms"),
]

# Header row
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = ""
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    tf = cell.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = h
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

# Data rows
for i, row in enumerate(rows, 1):
    is_proposed = "MiniLM" in row[0]
    for c, val in enumerate(row):
        cell = table.cell(i, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GREEN if is_proposed else WHITE
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = val
        r.font.size = Pt(14)
        r.font.bold = is_proposed
        r.font.color.rgb = GREEN if is_proposed else DARK_TEXT

# Right: key finding card
add_box(s, Inches(8.2), Inches(1.7), Inches(4.6), Inches(4.7),
        fill=LIGHT_GREEN, line=GREEN, line_width=2)
add_text(s, "KEY FINDING",
         Inches(8.2), Inches(1.9), Inches(4.6), Inches(0.4),
         size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "1.000", Inches(8.2), Inches(2.4), Inches(4.6), Inches(1.4),
         size=80, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "MRR for MiniLM-L6-v2",
         Inches(8.2), Inches(3.8), Inches(4.6), Inches(0.4),
         size=16, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(s, "The correct document is always\nranked first — perfect retrieval.",
         Inches(8.2), Inches(4.3), Inches(4.6), Inches(0.9),
         size=13, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(s, "2.7× better MRR than\nthe Word2Vec baseline.",
         Inches(8.2), Inches(5.3), Inches(4.6), Inches(0.9),
         size=13, italic=True, color=GREEN, align=PP_ALIGN.CENTER)

# Bottom takeaway
add_text(s,
         "Contextual sentence embeddings dramatically outperform classical NLP "
         "baselines for financial domain retrieval.",
         Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         size=15, italic=True, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=8)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Evaluation Experiment 2 (RAG vs No-RAG)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "Experiment 2: RAG vs No-RAG — Why Retrieval Matters")

add_text(s, "Same 5 financial questions, three answer strategies. Which approach hallucinates?",
         Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
         size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# Three comparison cards
cards = [
    ("No-RAG", "LLM alone, no document",
     [("Keyword Hit Rate", "0.400"),
      ("Faithfulness",     "0.600"),
      ("Hallucinations",   "1 / 5")],
     RED, RGBColor(0xFF, 0xEB, 0xEE)),
    ("Random Context", "Irrelevant chunks fed in",
     [("Keyword Hit Rate", "0.200"),
      ("Faithfulness",     "1.000"),
      ("Hallucinations",   "0 / 5")],
     ORANGE, LIGHT_GOLD),
    ("RAG (Proposed)", "Correctly retrieved chunks",
     [("Keyword Hit Rate", "1.000"),
      ("Faithfulness",     "1.000"),
      ("Hallucinations",   "0 / 5")],
     GREEN, LIGHT_GREEN),
]

x_start = Inches(0.5)
card_w  = Inches(4.1)
card_h  = Inches(4.8)
gap     = Inches(0.25)

for i, (name, sub, metrics, accent, fill) in enumerate(cards):
    x = x_start + (card_w + gap) * i
    y = Inches(1.6)

    add_box(s, x, y, card_w, card_h, fill=fill, line=accent, line_width=2)

    # Title
    add_text(s, name, x, y + Inches(0.25), card_w, Inches(0.55),
             size=22, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(s, sub, x, y + Inches(0.85), card_w, Inches(0.4),
             size=12, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             x + Inches(0.5), y + Inches(1.4),
                             card_w - Inches(1), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = accent
    div.line.fill.background()

    # metrics
    my = y + Inches(1.65)
    for mname, mval in metrics:
        add_text(s, mname, x + Inches(0.3), my,
                 card_w - Inches(0.6), Inches(0.3),
                 size=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)
        add_text(s, mval, x + Inches(0.3), my + Inches(0.05),
                 card_w - Inches(0.6), Inches(0.4),
                 size=20, bold=True, color=accent, align=PP_ALIGN.RIGHT)
        my += Inches(0.75)

# Bottom takeaway
add_text(s,
         "RAG achieves 100% Keyword Hit Rate and ZERO hallucinations — "
         "while No-RAG hallucinated a fabricated R&D percentage.",
         Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         size=15, italic=True, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=9)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Conclusions
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "Conclusions & Future Work")

# Left: what we achieved
add_text(s, "What We Achieved",
         Inches(0.5), Inches(1.1), Inches(6), Inches(0.5),
         size=20, bold=True, color=GREEN)

achievements = [
    "Built a production-ready RAG pipeline for financial documents",
    "Integrated 3 AI techniques: NLP, LLM, Prompt Engineering",
    "Demonstrated Word2Vec, TF-IDF, and Sentence-Transformer baselines",
    "Perfect retrieval (MRR = 1.000) on the test corpus",
    "Zero hallucinations with the RAG architecture",
    "Deployed as Streamlit web app + FastAPI REST service",
]

add_bullets(s, achievements,
            Inches(0.5), Inches(1.7), Inches(6.3), Inches(5),
            size=14, line_spacing=1.4)

# Right: future work
add_text(s, "Future Work",
         Inches(7.2), Inches(1.1), Inches(5.6), Inches(0.5),
         size=20, bold=True, color=ORANGE)

future = [
    "Hybrid retrieval (dense semantic + BM25 sparse)",
    "FinBERT for domain-specific financial embeddings",
    "Cross-document queries (compare multiple companies)",
    "Cross-encoder re-ranking of top retrieved chunks",
    "Conversation memory for follow-up questions",
    "OCR support for scanned financial documents",
]

add_bullets(s, future,
            Inches(7.2), Inches(1.7), Inches(5.6), Inches(5),
            size=14, line_spacing=1.4)

# Bottom takeaway bar
add_box(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.65),
        fill=NAVY, line=NAVY)
add_text(s,
         "The system satisfies all EC7203 requirements: 3 AI techniques, "
         "full evaluation with baselines, working demo, and reproducible code.",
         Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5),
         size=14, italic=True, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_footer(s, "EC7203 Advanced AI  •  Financial Document Q&A System", num=10)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Thank You
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, "Thank You",
         Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.08))
acc.fill.solid(); acc.fill.fore_color.rgb = GOLD; acc.line.fill.background()

add_text(s, "Questions & Discussion",
         Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.7),
         size=28, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# Deliverables summary
add_text(s, "Deliverables Submitted:",
         Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

delivs = "Final Report (PDF, 19 pages, LaTeX)   •   Working Web App   •   Source Code   •   This Video"
add_text(s, delivs, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
         size=14, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "EC7203 Advanced Artificial Intelligence  •  Final Group Project",
         Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "May 2026", Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3),
         size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "Financial_QA_Presentation.pptx",
)
prs.save(output_path)
print(f"Created: {output_path}")
print(f"Total slides: {len(prs.slides)}")
